# -*- coding: utf-8 -*-
"""Экспорт финального отчёта разбора в .docx (Фаза D bonus).

Конвертирует markdown-отчёт (final_report.markdown) в Word-документ.
Поддерживает минимальный markdown: заголовки (#/##/###), списки (- ),
жирный (**...**), горизонтальные линии (---). Этого достаточно для
отчётов, которые генерирует assemble.py.

Возвращает bytes (.docx) — endpoint отдаёт как файл. python-docx уже в
requirements (1.2.0).
"""
from __future__ import annotations

import io
import logging
import os
import re
import shutil
import subprocess
import tempfile
from typing import Any, Optional

logger = logging.getLogger(__name__)

_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


def _add_runs_with_bold(paragraph, text: str) -> None:
    """Добавить текст в параграф, разбирая **bold**."""
    pos = 0
    for m in _BOLD_RE.finditer(text):
        if m.start() > pos:
            paragraph.add_run(text[pos:m.start()])
        run = paragraph.add_run(m.group(1))
        run.bold = True
        pos = m.end()
    if pos < len(text):
        paragraph.add_run(text[pos:])


def _apply_letterhead(doc, letterhead_png: bytes) -> None:
    """Вставить фирменный бланк-картинку в верхний колонтитул (шапка на каждой
    странице). Ширина — по контенту страницы. Never-raise."""
    try:
        section = doc.sections[0]
        usable = section.page_width - section.left_margin - section.right_margin
        header = section.header
        header.is_linked_to_previous = False
        para = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        try:
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception:
            pass
        para.add_run().add_picture(io.BytesIO(letterhead_png), width=usable)
    except Exception:
        logger.warning("letterhead apply skipped", exc_info=True)


def markdown_to_docx_bytes(markdown: str, *, title: str = "Аналитический разбор",
                           letterhead_png: bytes | None = None) -> bytes:
    """Конвертировать markdown-отчёт в .docx, вернуть bytes.

    letterhead_png — фирменный бланк-картинка (PNG/JPG) в шапку каждой страницы;
    None → документ ровно как раньше."""
    from docx import Document  # lazy: тяжёлый импорт

    doc = Document()
    # Метаданные документа.
    try:
        doc.core_properties.title = title
    except Exception:
        logger.debug("docx core_properties set skipped", exc_info=True)

    if letterhead_png:
        _apply_letterhead(doc, letterhead_png)

    for raw_line in (markdown or "").split("\n"):
        line = raw_line.rstrip()
        if not line:
            continue

        # Горизонтальная линия → пустой параграф-разделитель.
        if line.strip() in ("---", "***", "___"):
            doc.add_paragraph()
            continue

        # Заголовки.
        if line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        # Списки.
        elif line.lstrip().startswith(("- ", "* ")):
            content = line.lstrip()[2:]
            p = doc.add_paragraph(style="List Bullet")
            _add_runs_with_bold(p, content)
        else:
            p = doc.add_paragraph()
            _add_runs_with_bold(p, line)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _soffice_bin() -> str:
    """Путь к LibreOffice для конвертации DOCX→PDF. Пусто → недоступно."""
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p
    return ""


_PDF_FONT_REG = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
    "C:/Windows/Fonts/arial.ttf", "C:/Windows/Fonts/segoeui.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
]
_PDF_FONT_BOLD = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
    "C:/Windows/Fonts/arialbd.ttf", "C:/Windows/Fonts/segoeuib.ttf",
    "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
]


def _first_existing(paths: list[str]) -> Optional[str]:
    for p in paths:
        if os.path.isfile(p):
            return p
    return None


def pdf_available() -> bool:
    """PDF доступен, если есть fpdf2 (переносимо) ИЛИ LibreOffice."""
    try:
        import fpdf
        return _first_existing(_PDF_FONT_REG) is not None
    except Exception:
        return bool(_soffice_bin())


def _markdown_to_pdf_fpdf(markdown: str, *, title: str,
                          letterhead_png: bytes | None) -> Optional[bytes]:
    """Переносимый рендер Markdown → PDF (fpdf2 + DejaVu, идеальная кириллица).
    Бланк-картинка ставится сверху первой страницы. Never-raise."""
    try:
        from fpdf import FPDF
    except Exception:
        return None
    reg = _first_existing(_PDF_FONT_REG)
    bold = _first_existing(_PDF_FONT_BOLD) or reg
    if not reg:
        return None
    try:
        pdf = FPDF(format="A4")
        pdf.set_auto_page_break(auto=True, margin=18)
        pdf.set_margins(20, 18, 20)
        pdf.add_page()
        pdf.add_font("Body", "", reg)
        pdf.add_font("Body", "B", bold)
        pdf.set_title(title or "")
        usable = pdf.w - pdf.l_margin - pdf.r_margin
        if letterhead_png:
            try:
                pdf.image(io.BytesIO(letterhead_png), x=pdf.l_margin, w=usable)
                pdf.ln(4)
            except Exception:
                logger.debug("pdf letterhead skipped", exc_info=True)
        for raw in (markdown or "").split("\n"):
            line = raw.rstrip()
            if not line:
                pdf.ln(3)
                continue
            if line.strip() in ("---", "***", "___"):
                y = pdf.get_y() + 2
                pdf.line(pdf.l_margin, y, pdf.w - pdf.r_margin, y)
                pdf.ln(4)
                continue
            if line.startswith("### "):
                pdf.set_font("Body", "B", 12); pdf.multi_cell(usable, 7, line[4:].strip())
            elif line.startswith("## "):
                pdf.set_font("Body", "B", 14); pdf.multi_cell(usable, 8, line[3:].strip())
            elif line.startswith("# "):
                pdf.set_font("Body", "B", 17); pdf.multi_cell(usable, 9, line[2:].strip())
            elif line.lstrip().startswith(("- ", "* ")):
                pdf.set_font("Body", "", 11)
                pdf.multi_cell(usable, 6, "•  " + line.lstrip()[2:], markdown=True)
            else:
                pdf.set_font("Body", "", 11)
                pdf.multi_cell(usable, 6, line, markdown=True)
        out = pdf.output()
        return bytes(out)
    except Exception:
        logger.warning("fpdf pdf render failed", exc_info=True)
        return None


def docx_bytes_to_pdf(docx_bytes: bytes) -> Optional[bytes]:
    """DOCX → PDF через LibreOffice headless (сохраняет бланк/шапку/стили).
    None, если LibreOffice недоступен или конвертация не удалась. Never-raise.
    Безопасно: subprocess со списком аргументов, без shell, во временной папке."""
    soffice = _soffice_bin()
    if not soffice or not docx_bytes:
        return None
    tmp = tempfile.mkdtemp(prefix="tb_pdf_")
    try:
        src = os.path.join(tmp, "doc.docx")
        with open(src, "wb") as f:
            f.write(docx_bytes)
        # LibreOffice кладёт doc.pdf рядом; свой профиль — чтобы не конфликтовать.
        proc = subprocess.run(
            [soffice, "--headless", "--nologo", "--nofirststartwizard",
             f"-env:UserInstallation=file://{os.path.join(tmp, 'profile')}",
             "--convert-to", "pdf", "--outdir", tmp, src],
            capture_output=True, timeout=120, check=False)
        out = os.path.join(tmp, "doc.pdf")
        if os.path.isfile(out):
            with open(out, "rb") as f:
                return f.read()
        logger.warning("docx→pdf: no output (%s)", (proc.stderr or b"")[:200])
        return None
    except Exception:
        logger.warning("docx→pdf conversion failed", exc_info=True)
        return None
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def markdown_to_pdf_bytes(markdown: str, *, title: str = "Документ",
                          letterhead_png: bytes | None = None) -> Optional[bytes]:
    """Markdown → PDF. Переносимый путь fpdf2 (кириллица через DejaVu, бланк
    сверху); если fpdf нет — фолбэк DOCX→PDF через LibreOffice. None — если ни то,
    ни другое недоступно (роут тогда отдаёт DOCX)."""
    pdf = _markdown_to_pdf_fpdf(markdown, title=title, letterhead_png=letterhead_png)
    if pdf:
        return pdf
    docx = markdown_to_docx_bytes(markdown, title=title, letterhead_png=letterhead_png)
    return docx_bytes_to_pdf(docx)


def report_to_docx_bytes(final_report: dict[str, Any], *, playbook_name: str = "") -> bytes:
    """Из final_report-объекта собрать .docx."""
    md = (final_report or {}).get("markdown", "") if isinstance(final_report, dict) else ""
    title = f"Разбор: {playbook_name}" if playbook_name else "Аналитический разбор"
    return markdown_to_docx_bytes(md, title=title)


# ──────────────────────────────────────────────────────────────────────
# XLSX (нативный Excel-отчёт) — openpyxl уже в requirements
# ──────────────────────────────────────────────────────────────────────

# Управляющие символы, запрещённые в OOXML/Excel (кроме \t \n \r): openpyxl
# бросает IllegalCharacterError. Чистим перед записью ячеек.
_XLSX_ILLEGAL_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def _norm_cell(v: Any) -> Any:
    """Значение ячейки: числа/строки как есть, None → '', остальное → str.
    Excel сам различит число и текст — числа не заворачиваем в строку.
    Управляющие символы вырезаются (иначе openpyxl падает IllegalCharacterError)."""
    if v is None:
        return ""
    if isinstance(v, bool):
        return "да" if v else "нет"
    if isinstance(v, (int, float)):
        return v
    if isinstance(v, str):
        s = _XLSX_ILLEGAL_RE.sub("", v).strip()
        # «1 234,50» / «1234.5» → число, чтобы Excel считал по колонке
        cleaned = s.replace(" ", "").replace(" ", "").replace(",", ".")
        if cleaned and re.fullmatch(r"-?\d+(\.\d+)?", cleaned):
            try:
                num = float(cleaned)
                return int(num) if num.is_integer() else num
            except (TypeError, ValueError):
                return s
        return s
    return str(v)


def rows_to_xlsx_bytes(sheets: list[dict]) -> Optional[bytes]:
    """Структурированные данные → нативный .xlsx (openpyxl).

    sheets: [{"name": str, "columns": [str], "rows": [[cell, ...], ...]}].
    Шапка жирная, строки замораживаются, ширина колонок — по содержимому.
    Числовые строки распознаются как числа (Excel считает суммы). None если
    openpyxl недоступен."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter
    except Exception:
        logger.warning("openpyxl недоступен — xlsx не собран", exc_info=True)
        return None
    if not sheets:
        sheets = [{"name": "Отчёт", "columns": [], "rows": []}]
    wb = Workbook()
    wb.remove(wb.active)
    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="2F5496")
    wrap = Alignment(vertical="top", wrap_text=True)
    for i, sh in enumerate(sheets):
        raw_name = str((sh or {}).get("name") or f"Лист {i + 1}").strip() or f"Лист {i + 1}"
        # Excel: имя листа ≤31 символа и без : \ / ? * [ ]
        name = re.sub(r"[:\\/?*\[\]]", " ", raw_name)[:31] or f"Лист {i + 1}"
        ws = wb.create_sheet(title=name)
        columns = [_XLSX_ILLEGAL_RE.sub("", str(c)) for c in ((sh or {}).get("columns") or [])]
        rows = (sh or {}).get("rows") or []
        widths: dict[int, int] = {}
        if columns:
            ws.append(columns)
            for ci, c in enumerate(columns, start=1):
                cell = ws.cell(row=1, column=ci)
                cell.font = head_font
                cell.fill = head_fill
                cell.alignment = wrap
                widths[ci] = len(str(c))
            ws.freeze_panes = "A2"
        for r in rows:
            cells = r if isinstance(r, (list, tuple)) else [r]
            ws.append([_norm_cell(v) for v in cells])
            for ci, v in enumerate(cells, start=1):
                widths[ci] = max(widths.get(ci, 8), min(60, len(str(v if v is not None else ""))))
        for ci, w in widths.items():
            ws.column_dimensions[get_column_letter(ci)].width = min(60, max(10, w + 2))
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ──────────────────────────────────────────────────────────────────────
# Markdown → красивый самодостаточный HTML (типографика + таблицы + печать)
# ──────────────────────────────────────────────────────────────────────

def _html_escape(s: str) -> str:
    return (s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


def _md_inline_html(s: str) -> str:
    """Инлайн-markdown → HTML: экранируем, потом **жирный** и `код`."""
    out = _html_escape(s)
    out = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", out)
    out = re.sub(r"`([^`]+)`", r"<code>\1</code>", out)
    return out


def markdown_to_html_bytes(markdown: str, *, title: str = "Документ",
                           accent: str = "#1f6feb") -> bytes:
    """Markdown → ОФОРМЛЕННЫЙ самодостаточный HTML (utf-8 bytes): типографика,
    стилизованные таблицы, печать в PDF (@media print), светлая тема. Всё инлайн,
    без внешних ресурсов — открывается в браузере и печатается как есть."""
    lines = (markdown or "").splitlines()
    body: list[str] = []
    i, n = 0, len(lines)
    para: list[str] = []

    def _flush_para() -> None:
        if para:
            body.append("<p>" + "<br>".join(_md_inline_html(x) for x in para) + "</p>")
            para.clear()

    while i < n:
        ln = lines[i].rstrip()
        s = ln.strip()
        h = re.match(r"^(#{1,6})\s+(.+)$", s)
        if h:
            _flush_para()
            lvl = min(3, len(h.group(1)))
            body.append(f"<h{lvl}>{_md_inline_html(h.group(2).strip())}</h{lvl}>")
            i += 1
            continue
        if s.startswith("|") and i + 1 < n and re.match(r"^\|?[\s:|-]+\|?$", lines[i + 1].strip()):
            _flush_para()
            cells = lambda r: [c.strip() for c in r.strip().strip("|").split("|")]
            cols = cells(s)
            rows = []
            j = i + 2
            while j < n and lines[j].strip().startswith("|"):
                rows.append(cells(lines[j]))
                j += 1
            thead = "".join(f"<th>{_md_inline_html(c)}</th>" for c in cols)
            trs = []
            for r in rows:
                tds = "".join(f"<td>{_md_inline_html(r[k] if k < len(r) else '')}</td>"
                              for k in range(len(cols)))
                trs.append(f"<tr>{tds}</tr>")
            body.append(f"<table><thead><tr>{thead}</tr></thead><tbody>{''.join(trs)}</tbody></table>")
            i = j
            continue
        if re.match(r"^[-*]\s+", s):
            _flush_para()
            items = []
            while i < n and re.match(r"^[-*]\s+", lines[i].strip()):
                items.append(re.sub(r"^[-*]\s+", "", lines[i].strip()))
                i += 1
            body.append("<ul>" + "".join(f"<li>{_md_inline_html(it)}</li>" for it in items) + "</ul>")
            continue
        if not s:
            _flush_para()
            i += 1
            continue
        if set(s) <= {"-", "=", "_"} and len(s) >= 3:
            _flush_para()
            body.append("<hr>")
            i += 1
            continue
        para.append(s)
        i += 1
    _flush_para()

    html = f"""<!doctype html><html lang="ru"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{_html_escape(title)}</title>
<style>
  :root{{--accent:{accent};--ink:#1f2733;--muted:#5a6472;--line:#e4e7ec;--soft:#eef3fb;}}
  *{{box-sizing:border-box}}
  body{{margin:0;background:#f4f6fa;color:var(--ink);
    font:16px/1.6 -apple-system,"Segoe UI",Roboto,Helvetica,Arial,sans-serif}}
  .doc{{max-width:860px;margin:24px auto;background:#fff;padding:48px 56px;
    border-radius:12px;box-shadow:0 2px 20px rgba(20,26,36,.06)}}
  .brandbar{{height:6px;background:var(--accent);border-radius:6px;margin:-20px 0 28px}}
  h1{{font-size:28px;letter-spacing:-.02em;margin:.2em 0 .4em}}
  h2{{font-size:15px;text-transform:uppercase;letter-spacing:.06em;color:var(--accent);
    margin:1.8em 0 .6em}}
  h3{{font-size:17px;margin:1.4em 0 .4em}}
  p{{margin:.5em 0}}
  ul{{margin:.4em 0 .8em;padding-left:22px}} li{{margin:.25em 0}}
  code{{background:var(--soft);padding:1px 5px;border-radius:4px;font-size:.9em}}
  hr{{border:0;border-top:1px solid var(--line);margin:1.6em 0}}
  table{{width:100%;border-collapse:collapse;margin:1em 0;font-variant-numeric:tabular-nums}}
  th{{background:var(--accent);color:#fff;text-align:left;padding:9px 12px;font-size:13px;
    text-transform:uppercase;letter-spacing:.03em}}
  td{{padding:9px 12px;border-bottom:1px solid var(--line)}}
  tbody tr:nth-child(even) td{{background:var(--soft)}}
  td:last-child,th:last-child{{text-align:right}}
  @media print{{body{{background:#fff}} .doc{{box-shadow:none;margin:0;border-radius:0;max-width:none}}}}
</style></head><body>
<div class="doc"><div class="brandbar"></div>
{''.join(body)}
</div></body></html>"""
    return html.encode("utf-8")


# ──────────────────────────────────────────────────────────────────────
# Markdown → таблицы (детерминированный парсер, без LLM) → XLSX
# ──────────────────────────────────────────────────────────────────────

def _markdown_tables_to_sheets(markdown: str) -> list[dict]:
    """Достаёт markdown-таблицы (| a | b |) в формат sheets для rows_to_xlsx_bytes.

    Каждая таблица → отдельный лист; заголовок берётся из ближайшей строки
    «## …» перед таблицей. Разделитель |---|---| пропускается. Без таблиц —
    один лист с построчным текстом (чтобы xlsx всегда собрался, а не падал)."""
    lines = (markdown or "").splitlines()
    sheets: list[dict] = []
    last_heading = "Данные"
    i, n = 0, len(lines)
    while i < n:
        ln = lines[i].strip()
        h = re.match(r"^#{1,6}\s+(.+)$", ln)
        if h:
            last_heading = h.group(1).strip()[:40]
            i += 1
            continue
        # начало таблицы: строка-заголовок с | и следующая — разделитель ---
        if ln.startswith("|") and i + 1 < n and re.match(r"^\|?[\s:|-]+\|?$", lines[i + 1].strip()):
            def _cells(row: str) -> list[str]:
                return [c.strip() for c in row.strip().strip("|").split("|")]
            columns = _cells(ln)
            rows: list[list[str]] = []
            j = i + 2
            while j < n and lines[j].strip().startswith("|"):
                rows.append(_cells(lines[j]))
                j += 1
            sheets.append({"name": last_heading or f"Лист {len(sheets)+1}",
                           "columns": columns, "rows": rows})
            i = j
            continue
        i += 1
    if not sheets:
        body = [[l] for l in lines if l.strip()][:500]
        sheets = [{"name": "Документ", "columns": ["Строка"], "rows": body}]
    return sheets


def markdown_to_xlsx_bytes(markdown: str, *, title: str = "Документ") -> Optional[bytes]:
    """Markdown-документ → .xlsx: таблицы из markdown становятся листами Excel
    (числа распознаются как числа). Детерминированно, без LLM. None если
    openpyxl недоступен."""
    return rows_to_xlsx_bytes(_markdown_tables_to_sheets(markdown))


# ──────────────────────────────────────────────────────────────────────
# Markdown → PPTX (python-pptx уже в requirements — используется парсером)
# ──────────────────────────────────────────────────────────────────────

def _md_to_pptx_sections(markdown: str, title: str) -> list[dict]:
    """Markdown → секции с УПОРЯДОЧЕННЫМИ элементами (буллеты и ТАБЛИЦЫ отдельно).

    section = {"title": str, "elements": [{"kind":"bullets","items":[...]},
    {"kind":"table","columns":[...],"rows":[[...]]}]}. Таблицы остаются
    таблицами (а не схлопываются в текст) — их рендерим настоящей pptx-таблицей."""
    def _clean(s: str) -> str:
        return _BOLD_RE.sub(r"\1", s).strip()

    def _cells(row: str) -> list[str]:
        return [_clean(c) for c in row.strip().strip("|").split("|")]

    lines = (markdown or "").splitlines()
    sections: list[dict] = []
    cur = {"title": title, "elements": []}

    def _bullet(items_src: list, txt: str) -> None:
        if not cur["elements"] or cur["elements"][-1]["kind"] != "bullets":
            cur["elements"].append({"kind": "bullets", "items": []})
        cur["elements"][-1]["items"].append(txt)

    i, n = 0, len(lines)
    while i < n:
        ln = lines[i].rstrip()
        h = re.match(r"^(#{1,6})\s+(.+)$", ln)
        if h:
            if cur["elements"] or cur["title"] != title:
                sections.append(cur)
            cur = {"title": _clean(h.group(2))[:90], "elements": []}
            i += 1
            continue
        s = ln.strip()
        # таблица: строка с | и следующая — разделитель |---|
        if s.startswith("|") and i + 1 < n and re.match(r"^\|?[\s:|-]+\|?$", lines[i + 1].strip()):
            columns = _cells(s)
            rows: list[list[str]] = []
            j = i + 2
            while j < n and lines[j].strip().startswith("|"):
                rows.append(_cells(lines[j]))
                j += 1
            cur["elements"].append({"kind": "table", "columns": columns, "rows": rows})
            i = j
            continue
        if not s or set(s) <= {"-", "*", "—", "="}:
            i += 1
            continue
        item = re.sub(r"^[-*]\s+", "", s).strip()
        if item:
            _bullet(cur["elements"], _clean(item)[:200])
        i += 1
    if cur["elements"] or cur["title"] != title:
        sections.append(cur)
    return sections


def markdown_to_pptx_bytes(markdown: str, *, title: str = "Документ",
                           accent: tuple[int, int, int] = (31, 111, 235),
                           subtitle: str = "") -> Optional[bytes]:
    """Markdown → ОФОРМЛЕННАЯ презентация .pptx (16:9): титульный слайд с
    акцентным фоном, секции с цветной шапкой, буллеты И НАСТОЯЩИЕ таблицы
    (прайс остаётся таблицей). Не «текст в pptx», а презентабельный деком.
    None если python-pptx недоступен. Реальная генерация, без заглушек."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except Exception:
        logger.warning("python-pptx недоступен — pptx не собран", exc_info=True)
        return None

    ACCENT = RGBColor(*accent)
    ACCENT_DK = RGBColor(max(0, accent[0] - 40), max(0, accent[1] - 40), max(0, accent[2] - 40))
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    INK = RGBColor(0x20, 0x27, 0x33)
    MUTE = RGBColor(0x8A, 0x93, 0xA0)
    ROW_ALT = RGBColor(0xF0, 0xF4, 0xFB)

    sections = _md_to_pptx_sections(markdown, title)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def _fill(shape, color) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _text(tf, text, *, size, color, bold=False, align=PP_ALIGN.LEFT):
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        return p

    # ── Титульный слайд: акцентный фон ──
    try:
        s = prs.slides.add_slide(blank)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        _fill(bg, ACCENT)
        box = s.shapes.add_textbox(Inches(0.9), Inches(2.6), SW - Inches(1.8), Inches(2.2))
        tf = box.text_frame
        _text(tf, title[:120], size=40, color=WHITE, bold=True)
        p2 = tf.add_paragraph()
        r = p2.add_run()
        r.text = subtitle or "Коммерческое предложение"
        r.font.size = Pt(20)
        r.font.color.rgb = RGBColor(0xE6, 0xEE, 0xFF)
    except Exception:
        logger.debug("pptx title slide skipped", exc_info=True)

    def _new_section_slide(sec_title: str):
        s = prs.slides.add_slide(blank)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
        _fill(bar, ACCENT)
        tb = s.shapes.add_textbox(Inches(0.7), Inches(0.18), SW - Inches(1.4), Inches(0.8))
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(tf, sec_title[:90], size=26, color=WHITE, bold=True)
        return s

    def _add_table(slide, columns, rows, top):
        rows = [r for r in rows if any(str(c).strip() for c in r)][:14]
        ncols = max(1, len(columns) or (len(rows[0]) if rows else 1))
        nrows = 1 + len(rows)
        left, width = Inches(0.7), SW - Inches(1.4)
        height = Inches(min(5.0, 0.45 * nrows))
        gtable = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
        # шапка
        for c in range(ncols):
            cell = gtable.cell(0, c)
            _fill_cell(cell, ACCENT_DK)
            _cell_text(cell, columns[c] if c < len(columns) else "", size=15, color=WHITE, bold=True)
        # строки
        for ri, row in enumerate(rows, start=1):
            for c in range(ncols):
                cell = gtable.cell(ri, c)
                if ri % 2 == 0:
                    _fill_cell(cell, ROW_ALT)
                else:
                    _fill_cell(cell, WHITE)
                val = row[c] if c < len(row) else ""
                _cell_text(cell, val, size=14, color=INK,
                           bold=False, align=(PP_ALIGN.RIGHT if c == ncols - 1 and ncols > 1 else PP_ALIGN.LEFT))
        return top + height + Inches(0.25)

    def _fill_cell(cell, color) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    def _cell_text(cell, text, *, size, color, bold=False, align=PP_ALIGN.LEFT) -> None:
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(text)[:120]
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    def _add_bullets(slide, items, top):
        items = [b for b in items if b][:11]
        if not items:
            return top
        h = Inches(min(5.0, 0.42 * len(items) + 0.2))
        box = slide.shapes.add_textbox(Inches(0.7), top, SW - Inches(1.4), h)
        tf = box.text_frame
        tf.word_wrap = True
        for k, b in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            dot = p.add_run()
            dot.text = "•  "
            dot.font.size = Pt(16)
            dot.font.color.rgb = ACCENT
            run = p.add_run()
            run.text = b
            run.font.size = Pt(16)
            run.font.color.rgb = INK
        return top + h + Inches(0.2)

    for sec in sections[:40]:
        try:
            slide = _new_section_slide(sec["title"])
            top = Inches(1.5)
            for el in sec["elements"][:8]:
                if top > Inches(6.7):
                    break
                if el["kind"] == "table" and (el.get("columns") or el.get("rows")):
                    top = _add_table(slide, el.get("columns") or [], el.get("rows") or [], top)
                elif el["kind"] == "bullets":
                    top = _add_bullets(slide, el.get("items") or [], top)
        except Exception:
            logger.debug("pptx section slide skipped", exc_info=True)
            continue

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
