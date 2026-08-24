import io
import logging
from typing import Tuple

# Try importing libraries
try:
    import pypdf
except ImportError:
    pypdf = None

# pdfplumber - лучше для сложных PDF
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

# OCR для PDF с изображениями
try:
    from pdf2image import convert_from_bytes
    pdf2image_available = True
except ImportError:
    pdf2image_available = False

try:
    import pytesseract
    pytesseract_available = True
except ImportError:
    pytesseract_available = False

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

logger = logging.getLogger(__name__)

class FileParser:
    """
    Сервис для извлечения текста из различных форматов файлов (PDF, DOCX, PPTX, XLSX, TXT).
    """

    @staticmethod
    def parse_file(content: bytes, filename: str) -> Tuple[str, str]:
        """
        Определяет тип файла и извлекает из него текст.

        Args:
            content: Содержимое файла в байтах
            filename: Имя файла для определения расширения

        Returns:
            Tuple[str, str]: (извлеченный_текст, тип_документа)
        """
        ext = filename.lower().split('.')[-1] if '.' in filename else ''

        if ext == 'pdf':
            logger.info(f"Parsing PDF file: {filename}")
            return FileParser._parse_pdf(content), 'pdf'
        elif ext in ['docx', 'doc']:
            logger.info(f"Parsing DOCX file: {filename}")
            return FileParser._parse_docx(content), 'docx'
        elif ext in ['pptx', 'ppt']:
            logger.info(f"Parsing PPTX file: {filename}")
            return FileParser._parse_pptx(content), 'pptx'
        elif ext in ['xlsx', 'xls']:
            logger.info(f"Parsing Excel file: {filename}")
            return FileParser._parse_excel(content), 'excel'
        elif ext == 'csv':
            logger.info(f"Parsing CSV file: {filename}")
            return FileParser._parse_csv(content), 'csv'
        elif ext == 'txt':
             return FileParser._parse_txt(content), 'text'
        else:
            # Нет/неизвестное расширение (напр. файл «callinsight преза» без
            # .pptx). Раньше бинарник читался КАК ТЕКСТ: pptx/pdf превращались в
            # мегабайты мусора, распухала запись и падал аплоад (Supabase
            # ReadTimeout на 24 МБ «текста»). Определяем формат по сигнатуре
            # (magic bytes) и парсим по-настоящему; не распознали — прежний
            # текстовый фолбэк. Только для ветки неизвестного расширения —
            # известные форматы выше не затрагиваются.
            sig = content[:4]
            if sig == b'%PDF':
                logger.info(f"{filename}: PDF определён по сигнатуре")
                return FileParser._parse_pdf(content), 'pdf'
            if sig == b'PK\x03\x04':  # OOXML/zip: pptx/docx/xlsx
                try:
                    import zipfile
                    names = zipfile.ZipFile(io.BytesIO(content)).namelist()
                    if any(n.startswith('ppt/') for n in names):
                        logger.info(f"{filename}: PPTX определён по содержимому zip")
                        return FileParser._parse_pptx(content), 'pptx'
                    if any(n.startswith('word/') for n in names):
                        logger.info(f"{filename}: DOCX определён по содержимому zip")
                        return FileParser._parse_docx(content), 'docx'
                    if any(n.startswith('xl/') for n in names):
                        logger.info(f"{filename}: XLSX определён по содержимому zip")
                        return FileParser._parse_excel(content), 'excel'
                except Exception as e:
                    logger.debug(f"zip-sniff не удался ({filename}): {e}")
            # По умолчанию пробуем как текст
            logger.info(f"Treating {filename} as text file")
            return FileParser._parse_txt(content), 'text'

    @staticmethod
    def _parse_csv(content: bytes) -> str:
        """CSV → аккуратная Markdown-таблица (было: сырой текст с запятыми).

        Структурированная таблица даёт LLM чёткую сетку колонок/строк — числа
        из финансовых csv извлекаются надёжнее, чем из плоского comma-текста.
        Разделитель определяется автоматически (, ; \\t). Большие файлы —
        первые 500 строк (остальное отражено в подписи)."""
        import csv
        import io as _io

        raw = FileParser._parse_txt(content)
        if not raw.strip():
            return ""
        # автоопределение разделителя по первой непустой строке
        sample = "\n".join(raw.splitlines()[:5])
        delimiter = ","
        try:
            delimiter = csv.Sniffer().sniff(sample, delimiters=",;\t").delimiter
        except Exception:
            for d in (";", "\t", ","):
                if d in sample:
                    delimiter = d
                    break
        try:
            reader = list(csv.reader(_io.StringIO(raw), delimiter=delimiter))
        except Exception:
            return raw  # фолбэк на сырой текст
        rows = [r for r in reader if any((c or "").strip() for c in r)]
        if not rows:
            return raw
        total = len(rows)
        rows = rows[:500]
        header = rows[0]
        ncol = max(len(r) for r in rows)
        header = header + [""] * (ncol - len(header))

        def _fmt(r):
            r = list(r) + [""] * (ncol - len(r))
            return "| " + " | ".join((c or "").strip().replace("|", "/") for c in r) + " |"

        md = [_fmt(header), "| " + " | ".join(["---"] * ncol) + " |"]
        md += [_fmt(r) for r in rows[1:]]
        out = "\n".join(md)
        if total > 500:
            out += f"\n\n_(показаны первые 500 из {total} строк)_"
        return out

    @staticmethod
    def _parse_txt(content: bytes) -> str:
        # Пытаемся декодировать как текст с разными кодировками
        try:
            return content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return content.decode('cp1251') # Common for Russian Windows
            except UnicodeDecodeError:
                try:
                    return content.decode('latin-1')
                except Exception:
                    logger.warning("Could not decode text file")
                    return ""

    @staticmethod
    def _parse_pdf(content: bytes) -> str:
        """
        Извлечение текста из PDF с несколькими fallback методами:
        1. pdfplumber (лучше для сложных PDF с таблицами и форматированием)
        2. pypdf (базовый метод)
        3. OCR (для PDF с изображениями/сканами)
        """
        text = []
        pages_without_text = []

        # Метод 1: pdfplumber (предпочтительный)
        if pdfplumber:
            try:
                with io.BytesIO(content) as f:
                    with pdfplumber.open(f) as pdf:
                        num_pages = len(pdf.pages)
                        logger.info(f"PDF has {num_pages} pages (using pdfplumber)")
                        for i, page in enumerate(pdf.pages):
                            # Извлекаем текст
                            page_text = page.extract_text()
                            if page_text and len(page_text.strip()) > 10:
                                text.append(page_text)
                            else:
                                # Пробуем извлечь текст из таблиц
                                tables = page.extract_tables()
                                if tables:
                                    for table in tables:
                                        for row in table:
                                            row_text = " | ".join([str(cell) if cell else "" for cell in row])
                                            text.append(row_text)
                                else:
                                    pages_without_text.append(i)
                                    logger.debug(f"No text on page {i+1}, will try OCR")

                        if text:
                            result = "\n".join(text)
                            logger.info(f"pdfplumber extracted {len(result)} chars from PDF")
                            # Если много страниц без текста, пробуем OCR
                            if len(pages_without_text) > num_pages / 2:
                                logger.info(f"⚠️ {len(pages_without_text)}/{num_pages} pages without text, trying OCR...")
                                ocr_text = FileParser._ocr_pdf(content)
                                if ocr_text and len(ocr_text) > len(result):
                                    logger.info(f"✅ OCR extracted more text: {len(ocr_text)} chars")
                                    return ocr_text.replace('\x00', '')
                            return result.replace('\x00', '')
            except Exception as e:
                logger.warning(f"pdfplumber failed: {e}, trying pypdf...")

        # Метод 2: pypdf (fallback)
        if pypdf:
            try:
                with io.BytesIO(content) as f:
                    reader = pypdf.PdfReader(f)
                    num_pages = len(reader.pages)
                    logger.info(f"PDF has {num_pages} pages (using pypdf)")
                    for i, page in enumerate(reader.pages):
                        extracted = page.extract_text()
                        if extracted and len(extracted.strip()) > 10:
                            text.append(extracted)
                        else:
                            pages_without_text.append(i)

                result = "\n".join(text)
                if result:
                    logger.info(f"pypdf extracted {len(result)} chars from PDF")
                    return result.replace('\x00', '')
            except Exception as e:
                logger.warning(f"pypdf failed: {e}")

        # Метод 3: OCR (для PDF с изображениями)
        if not text or len("\n".join(text).strip()) < 100:
            logger.info("📷 PDF appears to be image-based, trying OCR...")
            ocr_text = FileParser._ocr_pdf(content)
            if ocr_text:
                return ocr_text.replace('\x00', '')

        if not pypdf and not pdfplumber:
            logger.error("Neither pdfplumber nor pypdf library installed")

        return "\n".join(text).replace('\x00', '') if text else ""

    @staticmethod
    def _ocr_pdf(content: bytes) -> str:
        """
        OCR для PDF с изображениями (презентации, сканы).
        Требует: pdf2image, pytesseract, и установленный Tesseract OCR.
        """
        if not pdf2image_available:
            logger.warning("pdf2image not installed, OCR unavailable. Install: pip install pdf2image")
            return ""

        if not pytesseract_available:
            logger.warning("pytesseract not installed, OCR unavailable. Install: pip install pytesseract")
            return ""

        try:
            # Конвертируем PDF в изображения
            logger.info("🔄 Converting PDF to images for OCR...")
            images = convert_from_bytes(content, dpi=200)
            logger.info(f"📄 Converted {len(images)} pages to images")

            text_parts = []
            for i, image in enumerate(images):
                try:
                    # OCR с русским и английским языками
                    page_text = pytesseract.image_to_string(image, lang='rus+eng')
                    if page_text and len(page_text.strip()) > 10:
                        text_parts.append(f"--- Page {i+1} ---\n{page_text}")
                        logger.debug(f"OCR page {i+1}: {len(page_text)} chars")
                except Exception as e:
                    logger.warning(f"OCR failed for page {i+1}: {e}")

            result = "\n\n".join(text_parts)
            if result:
                logger.info(f"✅ OCR extracted {len(result)} chars from PDF")
            else:
                logger.warning("⚠️ OCR extracted no text from PDF")
            return result

        except Exception as e:
            logger.error(f"❌ OCR failed: {e}")
            # Проверяем типичные ошибки
            if "poppler" in str(e).lower():
                logger.error("💡 Hint: Install Poppler. Windows: choco install poppler / Linux: apt install poppler-utils")
            if "tesseract" in str(e).lower():
                logger.error("💡 Hint: Install Tesseract. Windows: choco install tesseract / Linux: apt install tesseract-ocr tesseract-ocr-rus")
            return ""

    @staticmethod
    def _parse_docx(content: bytes) -> str:
        if not Document:
            logger.error("python-docx library not installed")
            return ""

        try:
            with io.BytesIO(content) as f:
                doc = Document(f)
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)

                # Также извлекаем текст из таблиц
                for table in doc.tables:
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        full_text.append(" | ".join(row_text))

                return "\n".join(full_text)
        except Exception as e:
            logger.error(f"Error parsing DOCX: {e}")
            return ""

    @staticmethod
    def _parse_pptx(content: bytes) -> str:
        if not Presentation:
            logger.error("python-pptx library not installed")
            return ""

        try:
            with io.BytesIO(content) as f:
                prs = Presentation(f)
                text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    # Title
                    if slide.shapes.title:
                        slide_text.append(f"# {slide.shapes.title.text}")

                    # Other shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                            slide_text.append(shape.text)

                    if slide_text:
                        text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))

                return "\n\n".join(text)
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            return ""

    @staticmethod
    def _parse_excel(content: bytes) -> str:
        if not openpyxl:
            logger.error("openpyxl library not installed")
            return ""

        try:
            with io.BytesIO(content) as f:
                wb = openpyxl.load_workbook(f, data_only=True, read_only=True)
                text = []

                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    # Markdown-таблица с СОХРАНЕНИЕМ позиций колонок.
                    # Раньше None-ячейки выбрасывались и значения съезжали
                    # влево — числовые выгрузки (1С, сметы) приходили с
                    # перепутанными колонками; кроме того без |-рамки и
                    # разделителя таблицу не видели ни LLM, ни индукция
                    # методологий.
                    raw_rows = []
                    for row in ws.iter_rows(values_only=True):
                        cells = ["" if c is None else str(c).strip()
                                 for c in row]
                        while cells and not cells[-1]:
                            cells.pop()
                        if cells:
                            raw_rows.append(cells)
                        if len(raw_rows) >= 500:
                            break
                    if not raw_rows:
                        continue
                    width = max(len(r) for r in raw_rows)
                    lines = [f"--- Sheet: {sheet} ---"]
                    for i, r in enumerate(raw_rows):
                        padded = r + [""] * (width - len(r))
                        safe = [c.replace("|", "/").replace("\n", " ")
                                for c in padded]
                        lines.append("| " + " | ".join(safe) + " |")
                        if i == 0:
                            lines.append("|" + "---|" * width)
                    text.append("\n".join(lines))

                return "\n\n".join(text)
        except Exception as e:
            logger.error(f"Error parsing Excel: {e}")
            return ""
